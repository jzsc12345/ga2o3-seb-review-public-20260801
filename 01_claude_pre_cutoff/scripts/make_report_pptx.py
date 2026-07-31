#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
make_report_pptx.py -- 按 skills\\report-writer\\SKILL.md §5 的 outline.md 语法生成组会 PPTX。

体裁 A（推理文字为主）的排版规则内置：
  * 每页最多 2 图，图靠右占 40% 宽，文字占左侧主区（文字压倒图片）
  * `> ` 行进备注栏（deck 语法/求解器细节不进正文）
  * `---koupu: 概念名` 触发插入科普页骨架（定义/比喻/箭头链/为何在乎 四栏待填）
用法：
  python make_report_pptx.py outputs\\reports\\xxx_outline.md [-o outputs\\reports\\xxx.pptx]
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(r"D:\SILVACO_LOCAL")
W, H = Inches(13.333), Inches(7.5)          # 16:9
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)

KOUPU_FIELDS = ["大白话定义：", "比喻（括号体）：", "箭头因果链：", "本项目为什么在乎它："]


def parse_outline(path: Path) -> list[dict]:
    slides, cur = [], None
    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.rstrip()
        m_kp = re.match(r"^---koupu:\s*(.+)$", line)
        if line.startswith("# "):
            cur = {"title": line[2:].strip(), "body": [], "imgs": [], "notes": []}
            slides.append(cur)
        elif m_kp:
            slides.append({"title": f"概念扫盲：{m_kp.group(1).strip()}",
                           "body": KOUPU_FIELDS[:], "imgs": [], "notes":
                           ["科普页：填完四栏再交付；比喻必须括号括起来。"], "koupu": True})
            cur = None
        elif cur is None:
            continue
        elif line.startswith("> "):
            cur["notes"].append(line[2:])
        elif m := re.match(r"^!\[(.*?)\]\((.+?)\)\s*$", line):
            cur["imgs"].append((m.group(1), m.group(2)))
        elif line.strip():
            cur["body"].append(line)
    return slides


def build(slides: list[dict], out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    for i, s in enumerate(slides):
        sl = prs.slides.add_slide(blank)
        # 标题
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), W - Inches(1.0), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = s["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT if s.get("koupu") else DARK
        # 正文区：有图时占 58% 宽，无图占全宽 —— 文字永远是主角
        has_img = bool(s["imgs"])
        body_w = Inches(7.4) if has_img else W - Inches(1.0)
        bb = sl.shapes.add_textbox(Inches(0.5), Inches(1.3), body_w, H - Inches(1.9))
        tf = bb.text_frame
        tf.word_wrap = True
        for j, line in enumerate(s["body"]):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(16 if len(s["body"]) <= 8 else 14)
            para.space_after = Pt(6)
        # 图（≤2，右列纵排；路径不存在给灰框占位=「待补图」不硬凑）
        for k, (cap, ipath) in enumerate(s["imgs"][:2]):
            top = Inches(1.3) + k * Inches(2.9)
            ip = Path(ipath)
            if ip.exists():
                sl.shapes.add_picture(str(ip), Inches(8.1), top, width=Inches(4.7))
            else:
                ph = sl.shapes.add_textbox(Inches(8.1), top, Inches(4.7), Inches(2.6))
                ph.text_frame.text = f"[待补图]\n{cap}\n{ipath}"
                ph.text_frame.paragraphs[0].font.size = Pt(11)
            cb = sl.shapes.add_textbox(Inches(8.1), top + Inches(2.45), Inches(4.7), Inches(0.4))
            cp = cb.text_frame.paragraphs[0]
            cp.text = cap
            cp.font.size = Pt(11)
            cp.font.italic = True
        # 备注栏
        if s["notes"]:
            sl.notes_slide.notes_text_frame.text = "\n".join(s["notes"])
        # 页码
        nb = sl.shapes.add_textbox(W - Inches(0.9), H - Inches(0.45), Inches(0.6), Inches(0.3))
        nb.text_frame.text = str(i + 1)
        nb.text_frame.paragraphs[0].font.size = Pt(10)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"[pptx] {out}  ({len(slides)} 页)")
    missing = [(c, p) for s in slides for c, p in s["imgs"] if not Path(p).exists()]
    if missing:
        print(f"[待补图清单] {len(missing)} 张：")
        for c, p in missing:
            print(f"  - {c}: {p}")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("outline")
    ap.add_argument("-o", "--out", default=None)
    a = ap.parse_args()
    src = Path(a.outline)
    out = Path(a.out) if a.out else ROOT / "outputs" / "reports" / (src.stem.replace("_outline", "") + ".pptx")
    slides = parse_outline(src)
    if not slides:
        print("outline 里没有以『# 』开头的页标题", file=sys.stderr)
        return 2
    build(slides, out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
